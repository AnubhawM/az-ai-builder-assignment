#!/usr/bin/env python3
"""
Helper script for OpenClaw to generate PowerPoint presentations.
Usage: python create_pptx.py "Title" "slide1_title|slide1_content" "slide2_title|slide2_content" ...
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor

def create_presentation(title, slides_data, output_path="research_output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add title text box
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = 1  # Center
    
    # Content slides
    for slide_data in slides_data:
        if "|" in slide_data:
            slide_title, content = slide_data.split("|", 1)
        else:
            slide_title = slide_data
            content = ""
        
        content_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(content_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = slide_title.strip()
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Content
        if content:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for bullet in content.split("•"):
                if bullet.strip():
                    p = tf.add_paragraph()
                    p.text = "• " + bullet.strip()
                    p.font.size = Pt(18)
                    p.space_before = Pt(6)
    
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python create_pptx.py 'Title' 'slide1_title|content' 'slide2_title|content' ...")
        sys.exit(1)
    
    title = sys.argv[1]
    slides = sys.argv[2:]
    create_presentation(title, slides)
